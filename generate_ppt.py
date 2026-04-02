"""
PPT Generator — SkinVision AI Presentation
Generates a professional 16-slide PowerPoint presentation.
Run: python generate_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Constants ──
PRIMARY = RGBColor(10, 108, 158)       # #0A6C9E
PRIMARY_DARK = RGBColor(8, 85, 128)    # #085580
ACCENT = RGBColor(0, 196, 180)         # #00C4B4
BG_DARK = RGBColor(15, 23, 42)         # #0F172A
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(148, 163, 184)
TEXT_PRIMARY = RGBColor(31, 41, 55)
TEXT_SECONDARY = RGBColor(100, 116, 139)
SUCCESS = RGBColor(16, 185, 129)
WARNING = RGBColor(245, 158, 11)
ERROR = RGBColor(239, 68, 68)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def add_bg(slide, color=BG_DARK):
    """Fill slide background with solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, y=0, height=Inches(0.08)):
    """Add teal accent bar at top."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, y, SLIDE_WIDTH, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    """Add a text box with styling."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=LIGHT_GRAY, bullet_color=ACCENT):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(8)
        p.level = 0


def add_card(slide, left, top, width, height, title, body,
             title_color=PRIMARY, body_color=TEXT_SECONDARY):
    """Add a card-like shape with title and body."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(248, 250, 252)
    shape.line.color.rgb = RGBColor(226, 232, 240)
    shape.line.width = Pt(1)

    # Title
    add_text_box(slide, left + Inches(0.3), top + Inches(0.2),
                 width - Inches(0.6), Inches(0.5),
                 title, font_size=16, color=title_color, bold=True)
    # Body
    add_text_box(slide, left + Inches(0.3), top + Inches(0.7),
                 width - Inches(0.6), height - Inches(1),
                 body, font_size=13, color=body_color)


def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Use blank layout
    blank_layout = prs.slide_layouts[6]

    # ════════════════════════════════════════════════════════
    # SLIDE 1: TITLE
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide, BG_DARK)
    add_accent_bar(slide)

    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
                 "SKIN DISEASE DETECTION & CLASSIFICATION",
                 font_size=42, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(2.6), Inches(11), Inches(0.8),
                 "Using Convolutional Neural Networks",
                 font_size=28, color=ACCENT, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.6),
                 "A Virtual Dermatologist Powered by AI",
                 font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Divider
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.3), Inches(2.3), Inches(0.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

    add_text_box(slide, Inches(1), Inches(5), Inches(11), Inches(0.6),
                 "University Final Year Project",
                 font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(5.7), Inches(11), Inches(1),
                 "Domain: Healthcare / Dermatology / Computer Vision\n"
                 "Technology: Python | TensorFlow | Streamlit | EfficientNetB3",
                 font_size=14, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════
    # SLIDE 2: PROBLEM STATEMENT
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide, BG_DARK)
    add_accent_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "PROBLEM STATEMENT", font_size=36, color=WHITE, bold=True)

    problems = [
        "Skin diseases affect BILLIONS globally — yet access to dermatologists is extremely limited",
        "Only 1 dermatologist per 30,000 patients in developing countries",
        "Early detection of Melanoma = 98%+ survival rate; delayed = 23%",
        "Existing AI tools accept ANY image without validation — producing misleading results",
        "No visual explainability — doctors cannot trust a black-box prediction",
        "No downloadable reports — results disappear after the session",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5), problems, font_size=20)

    add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.8),
                 "GAP: Need a transparent, validated, accessible AI screening tool",
                 font_size=22, color=ACCENT, bold=True)

    # ════════════════════════════════════════════════════════
    # SLIDE 3: OBJECTIVES
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide, BG_DARK)
    add_accent_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "PROJECT OBJECTIVES", font_size=36, color=WHITE, bold=True)

    objectives = [
        "1.  Build a CNN model for classifying 7 types of skin diseases (HAM10000 dataset)",
        "2.  Achieve 80-92% accuracy using Transfer Learning with EfficientNetB3",
        "3.  Implement Grad-CAM for visual explainability (AI focus heatmaps)",
        "4.  Build 4-layer image validation (skin-tone, brightness, edge density)",
        "5.  Design a clinical-grade Streamlit web app with 7 interactive pages",
        "6.  Generate downloadable PDF diagnostic reports with disclaimers",
        "7.  Compare & improve upon existing open-source implementations",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5), objectives, font_size=20, color=WHITE)

    # ════════════════════════════════════════════════════════
    # SLIDE 4: DATASET — HAM10000
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide, BG_DARK)
    add_accent_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "DATASET: HAM10000", font_size=36, color=WHITE, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
                 "Human Against Machine with 10,000 training images  |  Harvard Dataverse / Kaggle",
                 font_size=16, color=LIGHT_GRAY)

    classes = [
        ("Melanocytic Nevi", "6,705", "66.9%", "Low"),
        ("Melanoma", "1,113", "11.1%", "CRITICAL"),
        ("Benign Keratosis", "1,099", "11.0%", "Low"),
        ("Basal Cell Carcinoma", "514", "5.1%", "High"),
        ("Actinic Keratosis", "327", "3.3%", "Moderate"),
        ("Vascular Lesions", "142", "1.4%", "Low"),
        ("Dermatofibroma", "115", "1.1%", "Low"),
    ]

    # Table header
    y_start = Inches(2.2)
    headers = ["Disease Class", "Images", "Percentage", "Severity"]
    x_positions = [Inches(1), Inches(5.5), Inches(7.5), Inches(9.5)]
    widths = [Inches(4), Inches(1.8), Inches(1.8), Inches(2)]

    for x, w, h in zip(x_positions, widths, headers):
        add_text_box(slide, x, y_start, w, Inches(0.4), h,
                     font_size=14, color=ACCENT, bold=True)

    for i, (name, count, pct, sev) in enumerate(classes):
        y = y_start + Inches(0.5) + Inches(i * 0.55)
        vals = [name, count, pct, sev]
        colors = [WHITE, WHITE, WHITE, SUCCESS if sev == "Low" else WARNING if sev == "Moderate" else ERROR]
        for x, w, v, c in zip(x_positions, widths, vals, colors):
            add_text_box(slide, x, y, w, Inches(0.4), v, font_size=14, color=c)

    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
                 "Key Challenge: Heavy class imbalance — solved with augmentation + class weights",
                 font_size=16, color=WARNING, bold=True)

    # ════════════════════════════════════════════════════════
    # SLIDE 5: METHODOLOGY
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide, BG_DARK)
    add_accent_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "METHODOLOGY", font_size=36, color=WHITE, bold=True)

    pipeline = [
        "1.  Data Collection — HAM10000 dataset (10,015 dermatoscopic images, 7 classes)",
        "2.  Preprocessing — Resize to 224x224, normalize [0,1], augmentation (flip, rotate, zoom)",
        "3.  Class Imbalance — Class weights + stratified split (80/10/10 train/val/test)",
        "4.  Model — EfficientNetB3 (ImageNet pre-trained) + custom classifier head",
        "5.  Training — Adam optimizer (lr=1e-4), categorical cross-entropy, early stopping",
        "6.  Evaluation — Accuracy, confusion matrix, classification report",
        "7.  Explainability — Grad-CAM heatmap generation",
        "8.  Deployment — Streamlit web app with validation, PDF reports, treatment info",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5), pipeline, font_size=18, color=WHITE)

    # ════════════════════════════════════════════════════════
    # SLIDE 6: MODEL ARCHITECTURE
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide, BG_DARK)
    add_accent_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "MODEL ARCHITECTURE", font_size=36, color=WHITE, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
                 "EfficientNetB3 — Transfer Learning from ImageNet",
                 font_size=20, color=ACCENT)

    arch = (
        "Input (224x224x3)  -->  EfficientNetB3 Base (FROZEN)  -->  GlobalAvgPooling2D\n\n"
        "-->  Dropout(0.3)  -->  Dense(256, ReLU)  -->  Dropout(0.3)  -->  Dense(7, Softmax)"
    )
    add_text_box(slide, Inches(0.8), Inches(2.2), Inches(11.5), Inches(1.5),
                 arch, font_size=20, color=WHITE)

    # Stats
    stats = [
        ("Total Parameters", "~12M"),
        ("Trainable Parameters", "~2M"),
        ("Base Model", "EfficientNetB3"),
        ("Pre-trained On", "ImageNet (1.4M images)"),
        ("Input Size", "224 x 224 x 3"),
        ("Output", "7 class probabilities (Softmax)"),
        ("Target Accuracy", "80-92%"),
    ]
    for i, (label, value) in enumerate(stats):
        y = Inches(4.0) + Inches(i * 0.45)
        add_text_box(slide, Inches(1.5), y, Inches(4), Inches(0.4),
                     label, font_size=14, color=LIGHT_GRAY)
        add_text_box(slide, Inches(5.5), y, Inches(5), Inches(0.4),
                     value, font_size=14, color=ACCENT, bold=True)

    # ════════════════════════════════════════════════════════
    # SLIDE 7: WHY EFFICIENTNET?
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide, BG_DARK)
    add_accent_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "WHY EfficientNetB3?", font_size=36, color=WHITE, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
                 "Compound Scaling = Best accuracy-to-parameter ratio",
                 font_size=18, color=ACCENT)

    models = [
        ("Model", "Parameters", "ImageNet Acc", "Choice"),
        ("VGG-16", "138M", "71.3%", "Too large, outdated"),
        ("ResNet-50", "25.6M", "76.0%", "Good but heavier"),
        ("EfficientNetB3", "12M", "81.6%", "BEST BALANCE"),
        ("EfficientNetB7", "66M", "84.3%", "Overkill for 10K images"),
    ]
    for i, (m, p, a, c) in enumerate(models):
        y = Inches(2.2) + Inches(i * 0.6)
        is_header = i == 0
        is_selected = i == 3
        mc = ACCENT if is_header else (SUCCESS if is_selected else WHITE)
        add_text_box(slide, Inches(1), y, Inches(3), Inches(0.5), m,
                     font_size=16 if not is_header else 14, color=mc, bold=is_header or is_selected)
        add_text_box(slide, Inches(4), y, Inches(2.5), Inches(0.5), p,
                     font_size=16 if not is_header else 14, color=mc, bold=is_header)
        add_text_box(slide, Inches(6.5), y, Inches(2.5), Inches(0.5), a,
                     font_size=16 if not is_header else 14, color=mc, bold=is_header)
        add_text_box(slide, Inches(9), y, Inches(3.5), Inches(0.5), c,
                     font_size=16 if not is_header else 14, color=mc, bold=is_header or is_selected)

    # ════════════════════════════════════════════════════════
    # SLIDE 8: GRAD-CAM
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide, BG_DARK)
    add_accent_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "GRAD-CAM EXPLAINABILITY", font_size=36, color=WHITE, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
                 "Gradient-weighted Class Activation Mapping — Making AI decisions visible",
                 font_size=18, color=ACCENT)

    steps = [
        "1.  Forward pass through the model with input image",
        "2.  Compute gradients of predicted class w.r.t. last convolutional layer",
        "3.  Global Average Pool gradients to get channel importance weights",
        "4.  Weighted sum of feature maps = raw activation heatmap",
        "5.  Apply ReLU (keep positive contributions only) + normalize",
        "6.  Resize to input dimensions and overlay with JET colormap",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4), steps, font_size=16, color=WHITE)

    add_text_box(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(1),
                 "WHY IT MATTERS: No existing GitHub implementation has Grad-CAM.\n"
                 "This is our KEY DIFFERENTIATOR — builds trust with clinicians.",
                 font_size=16, color=WARNING, bold=True)

    # ════════════════════════════════════════════════════════
    # SLIDE 9: IMAGE VALIDATION
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide, BG_DARK)
    add_accent_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "4-LAYER IMAGE VALIDATION", font_size=36, color=WHITE, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
                 "Inspired by DermaAI's is_skin() — enhanced with 3 additional checks",
                 font_size=16, color=ACCENT)

    checks = [
        ("Layer 1: File Validation", "Check file type (JPG/PNG), size (<10MB), image integrity"),
        ("Layer 2: Brightness Check", "Mean pixel brightness: reject <25 (too dark) or >245 (overexposed)"),
        ("Layer 3: Skin-Tone Detection", "HSV analysis [0,20,70]->[20,255,255] — DermaAI's exact range + extended"),
        ("Layer 4: Screenshot Detection", "Canny edge density >30% = UI screenshot with text content"),
    ]

    for i, (title, desc) in enumerate(checks):
        y = Inches(2.2) + Inches(i * 1.2)
        add_text_box(slide, Inches(1), y, Inches(5), Inches(0.5), title,
                     font_size=18, color=ACCENT, bold=True)
        add_text_box(slide, Inches(1), y + Inches(0.45), Inches(11), Inches(0.5), desc,
                     font_size=15, color=LIGHT_GRAY)

    # ════════════════════════════════════════════════════════
    # SLIDE 10: COMPARISON (KEY SLIDE)
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide, BG_DARK)
    add_accent_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "EXISTING SYSTEMS vs OUR SYSTEM", font_size=36, color=WHITE, bold=True)

    rows = [
        ("Feature", "Typical GitHub Projects", "SkinVision AI (Ours)"),
        ("Framework", "Flask (basic HTML/CSS)", "Streamlit (Modern, Interactive)"),
        ("Explainability", "NONE", "Grad-CAM Heatmaps"),
        ("Confidence", "Single output, no threshold", "Top-3 + 50% threshold + warnings"),
        ("Validation", "No input validation", "4-layer (HSV, brightness, edges)"),
        ("Reports", "Not available", "PDF download with branding"),
        ("Treatments", "Basic JSON (DermaAI)", "Detailed per-class recommendations"),
        ("UI Design", "Basic HTML forms", "Premium clinical CSS"),
        ("Ethics", "Minimal disclaimers", "Mandatory disclaimer gates"),
    ]

    for i, (feat, old, new) in enumerate(rows):
        y = Inches(1.5) + Inches(i * 0.6)
        is_header = i == 0
        fc = ACCENT if is_header else WHITE
        oc = ACCENT if is_header else ERROR
        nc = ACCENT if is_header else SUCCESS
        add_text_box(slide, Inches(0.5), y, Inches(3), Inches(0.5), feat,
                     font_size=14 if not is_header else 13, color=fc, bold=is_header)
        add_text_box(slide, Inches(3.5), y, Inches(4.5), Inches(0.5), old,
                     font_size=14 if not is_header else 13, color=oc, bold=is_header)
        add_text_box(slide, Inches(8), y, Inches(5), Inches(0.5), new,
                     font_size=14 if not is_header else 13, color=nc, bold=is_header)

    add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.5),
                 'Reference: DermaAI by FridahKimathi — github.com/FridahKimathi/Skin-Disease-Image-Classifier',
                 font_size=12, color=TEXT_SECONDARY)

    # ════════════════════════════════════════════════════════
    # SLIDE 11: WEB APPLICATION
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide, BG_DARK)
    add_accent_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "WEB APPLICATION — 7 PAGES", font_size=36, color=WHITE, bold=True)

    pages = [
        ("Home", "Hero section, impact stats, how-it-works steps"),
        ("About", "Project overview, objectives, dataset visualization"),
        ("How It Works", "Technical pipeline, validation checks, FAQ"),
        ("Model Details", "Architecture, 7 disease classes with severity"),
        ("Try the Detector", "Upload > Validate > Predict > Grad-CAM > PDF"),
        ("Comparison", "Feature table: Existing vs Our improvements"),
        ("Disclaimer & Ethics", "Medical disclaimer, bias, privacy, intended use"),
    ]

    for i, (name, desc) in enumerate(pages):
        y = Inches(1.6) + Inches(i * 0.75)
        add_text_box(slide, Inches(1), y, Inches(3.5), Inches(0.5),
                     name, font_size=18, color=ACCENT, bold=True)
        add_text_box(slide, Inches(4.5), y, Inches(8), Inches(0.5),
                     desc, font_size=16, color=LIGHT_GRAY)

    add_text_box(slide, Inches(0.8), Inches(7), Inches(11), Inches(0.4),
                 "Built with Streamlit + Custom CSS | Hospital-grade clinical design",
                 font_size=14, color=TEXT_SECONDARY)

    # ════════════════════════════════════════════════════════
    # SLIDE 12: DETECTOR DEMO FLOW
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide, BG_DARK)
    add_accent_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "DETECTOR FLOW — LIVE DEMO", font_size=36, color=WHITE, bold=True)

    flow = [
        "STEP 1:  User accepts medical disclaimer (mandatory gate)",
        "STEP 2:  Upload dermatoscopic image (JPG/PNG, max 10MB)",
        "STEP 3:  4-layer validation (brightness, skin-tone, edge density)",
        "STEP 4:  EfficientNetB3 CNN runs inference on preprocessed image",
        "STEP 5:  Grad-CAM heatmap generated (last convolutional layer)",
        "STEP 6:  Results displayed: predicted class + confidence + top-3",
        "STEP 7:  Treatment recommendations shown per disease class",
        "STEP 8:  Download PDF report with all findings + disclaimer",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5), flow, font_size=18, color=WHITE)

    # ════════════════════════════════════════════════════════
    # SLIDE 13: RESULTS & FEATURES
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide, BG_DARK)
    add_accent_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "RESULTS & FEATURES DELIVERED", font_size=36, color=WHITE, bold=True)

    features = [
        "7-class skin disease classification (HAM10000)",
        "80-92% target accuracy with EfficientNetB3",
        "Grad-CAM heatmap visualization (explainable AI)",
        "Top-3 predictions with confidence bars",
        "4-layer image validation (HSV skin detection, brightness, edges)",
        "50% confidence threshold with 'Inconclusive Result' warnings",
        "Treatment recommendations per disease class",
        "PDF diagnostic report with branding and disclaimers",
        "7-page clinical web application (Streamlit)",
        "Demo mode for presentations (image-aware, deterministic)",
    ]

    for i, feat in enumerate(features):
        y = Inches(1.5) + Inches(i * 0.55)
        add_text_box(slide, Inches(1.2), y, Inches(1), Inches(0.4),
                     "OK", font_size=14, color=SUCCESS, bold=True)
        add_text_box(slide, Inches(2), y, Inches(10), Inches(0.4),
                     feat, font_size=16, color=WHITE)

    # ════════════════════════════════════════════════════════
    # SLIDE 14: ETHICAL CONSIDERATIONS
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide, BG_DARK)
    add_accent_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "ETHICAL CONSIDERATIONS", font_size=36, color=WHITE, bold=True)

    ethics = [
        "NOT a diagnostic tool — educational and screening purposes only",
        "Dataset bias — HAM10000 may underrepresent darker skin tones",
        "Grad-CAM transparency — visual explanations of every prediction",
        "Mandatory disclaimer gates — users must acknowledge before using",
        "No data storage — all images processed in-session only",
        "Confidence thresholds — uncertain predictions are explicitly flagged",
        "Human oversight required — always recommends professional consultation",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5), ethics, font_size=18, color=WHITE)

    # ════════════════════════════════════════════════════════
    # SLIDE 15: LIMITATIONS & FUTURE WORK
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide, BG_DARK)
    add_accent_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "LIMITATIONS & FUTURE WORK", font_size=36, color=WHITE, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.5),
                 "LIMITATIONS", font_size=22, color=ERROR, bold=True)

    limitations = [
        "Dataset bias (skin tone underrepresentation)",
        "Class imbalance (115 vs 6,705 images)",
        "Not validated in clinical settings",
        "Only 7 disease categories",
        "Performance varies on non-dermoscopic images",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(2), Inches(5.5), Inches(4), limitations, font_size=16, color=LIGHT_GRAY)

    add_text_box(slide, Inches(7), Inches(1.4), Inches(5.5), Inches(0.5),
                 "FUTURE WORK", font_size=22, color=SUCCESS, bold=True)

    future = [
        "Diverse datasets (all skin types)",
        "Mobile app (TensorFlow Lite)",
        "Clinical validation with dermatologists",
        "Ensemble models for higher accuracy",
        "Federated learning across hospitals",
    ]
    add_bullet_list(slide, Inches(7), Inches(2), Inches(5.5), Inches(4), future, font_size=16, color=LIGHT_GRAY)

    # ════════════════════════════════════════════════════════
    # SLIDE 16: Q&A
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide, BG_DARK)
    add_accent_bar(slide)

    add_text_box(slide, Inches(1), Inches(2), Inches(11), Inches(1.2),
                 "THANK YOU", font_size=56, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Divider
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.3), Inches(2.3), Inches(0.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

    add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.8),
                 "Questions & Answers", font_size=28, color=ACCENT, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(5), Inches(11), Inches(1.5),
                 "SkinVision AI — Virtual Dermatologist\n"
                 "Skin Disease Detection & Classification Using CNN\n"
                 "University Final Year Project",
                 font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # ── Save ──
    output_path = os.path.join(os.path.dirname(__file__), 'SkinVisionAI_Presentation.pptx')
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
    return output_path


if __name__ == '__main__':
    create_presentation()
